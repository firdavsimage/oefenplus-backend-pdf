from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import landscape, A4

def convert_pptx_to_pdf(input_path, output_path):
    prs = Presentation(input_path)
    c = canvas.Canvas(output_path, pagesize=landscape(A4))
    width, height = landscape(A4)
    for slide in prs.slides:
        y = height - 50
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for line in shape.text.splitlines():
                    c.drawString(50, y, line)
                    y -= 15
        c.showPage()
    c.save()
