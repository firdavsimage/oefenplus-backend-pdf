import os
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from reportlab.pdfgen import canvas
from PIL import Image

def convert_docx_to_pdf(input_path, output_path):
    doc = Document(input_path)
    temp_txt = input_path + ".txt"
    with open(temp_txt, "w", encoding="utf-8") as f:
        for para in doc.paragraphs:
            f.write(para.text + "\n")
    text_to_pdf(temp_txt, output_path)
    os.remove(temp_txt)

def convert_pptx_to_pdf(input_path, output_path):
    pres = Presentation(input_path)
    temp_txt = input_path + ".txt"
    with open(temp_txt, "w", encoding="utf-8") as f:
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    f.write(shape.text + "\n")
    text_to_pdf(temp_txt, output_path)
    os.remove(temp_txt)

def convert_xlsx_to_pdf(input_path, output_path):
    wb = load_workbook(input_path)
    temp_txt = input_path + ".txt"
    with open(temp_txt, "w", encoding="utf-8") as f:
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                row_data = [str(cell.value or "") for cell in row]
                f.write("\t".join(row_data) + "\n")
    text_to_pdf(temp_txt, output_path)
    os.remove(temp_txt)

def image_to_pdf(input_path, output_path):
    image = Image.open(input_path)
    image.convert("RGB").save(output_path)

def text_to_pdf(txt_path, pdf_path):
    c = canvas.Canvas(pdf_path)
    with open(txt_path, "r", encoding="utf-8") as f:
        y = 800
        for line in f:
            c.drawString(50, y, line.strip())
            y -= 15
            if y < 50:
                c.showPage()
                y = 800
    c.save()

def convert_file(input_path, output_folder):
    ext = os.path.splitext(input_path)[1].lower()
    base = os.path.basename(input_path)
    filename = os.path.splitext(base)[0] + ".pdf"
    output_path = os.path.join(output_folder, filename)

    try:
        if ext == ".docx":
            convert_docx_to_pdf(input_path, output_path)
        elif ext == ".pptx":
            convert_pptx_to_pdf(input_path, output_path)
        elif ext == ".xlsx":
            convert_xlsx_to_pdf(input_path, output_path)
        elif ext in [".jpg", ".jpeg", ".png"]:
            image_to_pdf(input_path, output_path)
        else:
            return None, None
        return output_path, filename
    except Exception as e:
        print(f"[XATO] Konvertatsiyada xato: {e}")
        return None, None
